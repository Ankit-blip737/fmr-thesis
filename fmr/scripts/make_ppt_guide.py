import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(15, 23, 42)
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.add_paragraph()
    p2.text = subtitle_text
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(51, 65, 85)
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.3), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(59, 130, 246)
    line.line.color.rgb = RGBColor(59, 130, 246)
    
    return slide

def add_content_slide(prs, title, content_sections):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(15, 23, 42)
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(226, 232, 240)
    line.line.color.rgb = RGBColor(226, 232, 240)
    
    y_offset = 1.5
    for section in content_sections:
        if isinstance(section, dict) and "box" in section:
            # Estimate height based on text length
            text_len = len(section.get("text", ""))
            box_height = 1.0 + (text_len / 120.0) * 0.3
            
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_offset), Inches(9), Inches(box_height))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(219, 234, 254)
            box.line.color.rgb = RGBColor(191, 219, 254)
            
            tf_box = box.text_frame
            tf_box.word_wrap = True
            tf_box.margin_left = Inches(0.2)
            tf_box.margin_right = Inches(0.2)
            
            if section.get("title"):
                p_bt = tf_box.add_paragraph()
                p_bt.text = section["title"]
                p_bt.font.bold = True
                p_bt.font.size = Pt(18)
                p_bt.font.color.rgb = RGBColor(15, 23, 42)
            
            p_bc = tf_box.add_paragraph()
            p_bc.text = section["text"]
            p_bc.font.size = Pt(14)
            p_bc.font.color.rgb = RGBColor(30, 41, 59)
            
            y_offset += box_height + 0.2
            
        else:
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), Inches(9), Inches(0.5))
            tf_t = tx.text_frame
            tf_t.word_wrap = True
            
            if "title" in section:
                pt = tf_t.add_paragraph()
                pt.text = section["title"]
                pt.font.bold = True
                pt.font.size = Pt(20)
                pt.font.color.rgb = RGBColor(15, 23, 42)
                
            for bullet in section.get("bullets", []):
                pb = tf_t.add_paragraph()
                pb.text = "• " + bullet
                pb.font.size = Pt(16)
                pb.font.color.rgb = RGBColor(51, 65, 85)
                pb.level = 0
            
            y_offset += 0.6 + (len(section.get("bullets", [])) * 0.3)

def create_presentation():
    prs = Presentation()
    
    # 1. Title
    add_title_slide(prs, 
        "Faithful Medical Reasoning (FMR)", 
        "A Deep-Dive into the Safety Auditor for Medical AI\nAnkit Kumar · B.Tech Thesis")
    
    # 2. Module 1: The Problem
    add_content_slide(prs, "Module 1: The Clinical Problem", [
        {"title": "Why Does This Project Exist?", "bullets": [
            "Doctors are overwhelmed, and AI tools are being built to read X-rays and answer clinical questions automatically.",
            "The Danger: AI sometimes lies confidently (Hallucination). Worse, it might give the right answer for the wrong reason (guessing from text, not the image).",
            "Our Solution: We built a safety auditor that watches the AI think and pulls the emergency brake if it's not trustworthy."
        ]},
        {"box": True, "title": "The Novel Contribution", "text": "We measure 'grounding decay' step-by-step, fuse multiple faithfulness signals, and gate the output with a distribution-free safety guarantee."}
    ])
    
    # 3. Module 2: Key Terminology
    add_content_slide(prs, "Module 2: Key Terminology", [
        {"box": True, "title": "Vision-Language Model (VLM)", "text": "An AI (like MedVLM-R1) that can look at an image AND read text, generating an answer token-by-token. We use one that reasons step-by-step."},
        {"box": True, "title": "Grounding vs. Hallucination", "text": "Grounding means the AI's answer is based directly on the image pixels. Hallucination means the AI makes up an answer from text patterns."},
        {"box": True, "title": "Grounding Decay", "text": "Our central hypothesis: As a VLM generates more reasoning steps (longer chain-of-thought), it progressively 'forgets' the original image."}
    ])
    
    # 4. Module 3 & 9: Datasets
    add_content_slide(prs, "Module 3: Mock vs. Real Datasets", [
        {"title": "Why Two Types of Data?", "bullets": [
            "Mock (Synthetic): 300 computer-generated images. We control the ground truth and bounding boxes. Used to PROVE our mathematical formulas work.",
            "Real Datasets (VQA-RAD, PathVQA, SLAKE): Real X-rays and pathology slides. Used to PROVE the system applies to real hospital data."
        ]},
        {"box": True, "title": "The Bounding Box Limitation", "text": "Real datasets lack bounding boxes for diseases. Our Mock dataset provides them, allowing us to fully test Signal B (Spatial Grounding)."}
    ])
    
    # 5. Module 4: The 5-Stage Pipeline
    add_content_slide(prs, "Module 4: The 5-Stage Pipeline", [
        {"title": "How the System Processes Data", "bullets": [
            "Stage 1: Baselines — Establishes raw accuracy of models (usually ~25% on real data).",
            "Stage 2: Blind Test — Feeds blank images to expose hallucination (the 'Blind Gap').",
            "Stage 3: FMR Score — Computes trustworthiness (Signals A, B, C).",
            "Stage 4: Conformal Gate — Calculates the safety threshold for Answer/Abstain.",
            "Stage 5: Correction — Uses MedGemma as a second opinion (Future Work)."
        ]}
    ])
    
    # 6. Module 4 (cont): The Three Signals
    add_content_slide(prs, "Module 4: The Faithfulness Signals", [
        {"box": True, "title": "Signal A (Counterfactual / Image Reliance)", "text": "Show the model the real image, a blank image, and a mismatched image. If the answer changes, it relies on the image (Faithful)."},
        {"box": True, "title": "Signal B (Attention / Spatial Grounding)", "text": "Extract the model's internal attention heatmaps. Measure the overlap (IoU) with the actual disease bounding box."},
        {"box": True, "title": "Signal C (Consistency / Answer Stability)", "text": "Ask the exact same question 5 times with slight temperature randomness. A faithful model gives the same answer every time."},
        {"title": "Fusion", "bullets": ["FS = (w_A * A) + (w_B * B) + (w_C * C)"]}
    ])
    
    # 7. Module 7: Key Mathematical Concepts
    add_content_slide(prs, "Module 7: Key Mathematical Concepts", [
        {"title": "The Math Behind the Safety", "bullets": [
            "AUROC: Measures how well our FS score separates correct answers from incorrect ones. Perfect = 1.0.",
            "AURC (Risk-Coverage): As we answer more cases (Coverage), how does the error rate (Risk) grow? Lower is better."
        ]},
        {"box": True, "title": "Conformal Prediction (α, δ, τ)", "text": "α (Alpha) = Max tolerated error (e.g., 15%).\nδ (Delta) = Max chance the guarantee fails (e.g., 5%).\nτ (Tau) = The calculated FS threshold. If FS ≥ τ, we ANSWER. If FS < τ, we ABSTAIN."}
    ])
    
    # 8. Module 8: End-to-End Flow
    add_content_slide(prs, "Module 8: The End-to-End Flow", [
        {"title": "Putting It All Together", "bullets": [
            "1. Load medical images and questions.",
            "2. Run model to get answers and blind-test results.",
            "3. Compute Signals A, B, and C (5 passes = heavy GPU lifting).",
            "4. Fuse signals into the FS score.",
            "5. Split data into Calibration and Test sets.",
            "6. Find conformal threshold τ on calibration data.",
            "7. Output ANSWER or ABSTAIN for test data based on τ."
        ]}
    ])
    
    # 9. Module 5: The Dashboard
    add_content_slide(prs, "Module 5: The Interactive Dashboard", [
        {"title": "Visualizing the Pipeline Outputs", "bullets": [
            "Overview: High-level metrics, Baseline accuracy, and Replication verdicts.",
            "Diagnosis: Per-Step Grounding curves showing actual Grounding Decay.",
            "Measurement: AUROC and Risk-Coverage charts.",
            "Robustness: Ablation studies (Noise, Crop, Blur) proving the score works.",
            "Case Explorer: Searchable database of every single AI thought process."
        ]}
    ])
    
    # 10. Module 10: The Thesis Argument
    add_content_slide(prs, "Module 10: The Thesis Argument", [
        {"title": "What We Proved to the Committee", "bullets": [
            "1. Medical VLMs hallucinate (Proven by the Blind Test).",
            "2. Chain-of-thought doesn't fix it (Proven by Grounding Decay curves).",
            "3. We can detect it (Proven by high AUROC of our 3-signal FMR score).",
            "4. We can prevent it (Proven by the Conformal Abstention safety gate).",
            "5. It works on real medical data (Proven on VQA-RAD, PathVQA, and SLAKE)."
        ]},
        {"box": True, "title": "Conclusion", "text": "We built a rigorous, mathematically guaranteed safety auditor that is exactly what the medical AI industry needs."}
    ])
    
    prs.save('fmr/docs/FMR_Deep_Dive_Presentation.pptx')
    print("Saved to fmr/docs/FMR_Deep_Dive_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
