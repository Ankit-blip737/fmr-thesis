import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    
    # Add title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(15, 23, 42)
    
    # Add subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.add_paragraph()
    p2.text = subtitle_text
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(51, 65, 85)
    
    # Add a thin blue line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.3), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(59, 130, 246)
    line.line.color.rgb = RGBColor(59, 130, 246)
    
    return slide

def add_content_slide(prs, title, content_sections):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    
    # Add slide title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(15, 23, 42)
    
    # Add line below title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(226, 232, 240)
    line.line.color.rgb = RGBColor(226, 232, 240)
    
    # Content areas
    y_offset = 1.7
    for section in content_sections:
        if isinstance(section, dict) and "box" in section:
            # Draw a light blue box with text
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_offset), Inches(9), Inches(1.2))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(219, 234, 254) # light blue
            box.line.color.rgb = RGBColor(191, 219, 254)
            
            tf_box = box.text_frame
            tf_box.word_wrap = True
            tf_box.margin_left = Inches(0.2)
            tf_box.margin_right = Inches(0.2)
            
            if section.get("title"):
                p_bt = tf_box.add_paragraph()
                p_bt.text = section["title"]
                p_bt.font.bold = True
                p_bt.font.size = Pt(20)
                p_bt.font.color.rgb = RGBColor(15, 23, 42)
            
            p_bc = tf_box.add_paragraph()
            p_bc.text = section["text"]
            p_bc.font.size = Pt(16)
            p_bc.font.color.rgb = RGBColor(30, 41, 59)
            
            y_offset += 1.4
            
        else:
            # Normal text
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), Inches(9), Inches(0.5))
            tf_t = tx.text_frame
            tf_t.word_wrap = True
            
            if "title" in section:
                pt = tf_t.add_paragraph()
                pt.text = section["title"]
                pt.font.bold = True
                pt.font.size = Pt(22)
                pt.font.color.rgb = RGBColor(15, 23, 42)
                
            for bullet in section.get("bullets", []):
                pb = tf_t.add_paragraph()
                pb.text = "• " + bullet
                pb.font.size = Pt(18)
                pb.font.color.rgb = RGBColor(51, 65, 85)
                pb.level = 0
            
            y_offset += 0.8 + (len(section.get("bullets", [])) * 0.3)

def create_presentation():
    prs = Presentation()
    
    # 1. Title Slide
    add_title_slide(prs, 
        "Faithful Medical Reasoning (FMR)", 
        "A Safety Framework for Medical AI\nAnkit Kumar · B.Tech Thesis · Department of Computer Science")
    
    # 2. Clinical Problem
    add_content_slide(prs, "The Clinical Problem: Confident but Wrong", [
        {"title": "The Danger of Medical AI Hallucination", "bullets": [
            "Vision-Language Models (VLMs) are increasingly used to interpret medical images.",
            "Hallucination: Models often generate confident diagnoses that are entirely incorrect.",
            "The 'Clever Hans' Effect: Models can answer correctly by memorizing text patterns in the question, without looking at the image."
        ]},
        {"box": True, "title": "Grounding Decay", "text": "As models use 'Chain-of-Thought' reasoning, they progressively drift away from the visual evidence and rely solely on language patterns."}
    ])
    
    # 3. Innovation
    add_content_slide(prs, "Our Innovation: The FMR Framework", [
        {"title": "A First-of-its-Kind Safety Auditor", "bullets": [
            "We built FMR — a comprehensive system that audits the AI's thought process.",
            "Instead of blindly increasing accuracy, FMR:"
        ]},
        {"box": True, "title": "1. Measure", "text": "Quantifies how faithfully a VLM uses the medical image via a Faithfulness Score (0-1)."},
        {"box": True, "title": "2. Detect", "text": "Identifies grounding decay (attention drifting away from the image)."},
        {"box": True, "title": "3. Gate", "text": "Issues a mathematically guaranteed decision: ANSWER (safe) or ABSTAIN (defer)."}
    ])
    
    # 4. Architecture
    add_content_slide(prs, "System Architecture & Pipeline", [
        {"title": "End-to-End Evaluation Infrastructure", "bullets": [
            "Stage 1: Baselines — Raw accuracy of reasoning vs. non-reasoning models.",
            "Stage 2: Blind Test — Evaluates the model on blank images to detect pure hallucination.",
            "Stage 3: FMR Score — Computes our novel 3-signal faithfulness metric.",
            "Stage 4: Safety Gate — Applies Conformal Prediction to find the safety threshold.",
            "Stage 5: Live API — A FastAPI backend allowing real-time, dynamic inference."
        ]}
    ])
    
    # 5. Blind Test
    add_content_slide(prs, "Exposing Hallucination: The Blind Test", [
        {"title": "Proving When the Model Ignores the Image", "bullets": [
            "We ask the AI the same clinical question twice:",
            "  1. Once with the real X-ray.",
            "  2. Once with a completely blank (black) image."
        ]},
        {"box": True, "title": "The Blind Gap Metric", "text": "If the model gives the exact same correct answer on the blank image, it proves the model is guessing from text priors, not diagnosing from visual evidence."}
    ])
    
    # 6. FS Score
    add_content_slide(prs, "The Faithfulness Score (FS)", [
        {"title": "Fusing Three Independent Signals", "bullets": [
            "Signal A (Image Reliance): We swap the image with counterfactuals (blank, mismatched). A faithful model must change its answer.",
            "Signal B (Spatial Grounding): We extract the AI's internal attention heatmaps and calculate overlap (IoU) with the disease location.",
            "Signal C (Answer Consistency): We query the model 5 times with varying temperature. Grounded models answer consistently."
        ]},
        {"box": True, "title": "Fusion Formula", "text": "FS = w_A * A + w_B * B + w_C * C (A single score summarizing trustworthiness)"}
    ])
    
    # 7. Grounding Decay
    add_content_slide(prs, "Proving Grounding Decay", [
        {"title": "More Reasoning != More Faithful", "bullets": [
            "By tracking Signal B across the 'Chain-of-Thought', we proved:",
            "  • Step 1: Model focuses correctly on the disease region.",
            "  • Steps 2-3: Model's attention begins to drift.",
            "  • Steps 4+: Model essentially stops looking at the image and writes from memory."
        ]},
        {"box": True, "title": "Core Finding", "text": "Our system successfully graphs this decay, proving that forcing an AI to 'think more' can actually increase hallucination."}
    ])
    
    # 8. Conformal Gate
    add_content_slide(prs, "The Conformal Safety Gate", [
        {"title": "Distribution-Free Mathematical Guarantees", "bullets": [
            "Calibration: Split the data and define risk tolerance (e.g., max 15% error).",
            "Thresholding: The math finds a precise threshold (τ) that guarantees this error rate.",
            "The Guarantee: With 95% confidence, the system will maintain this safety level on unseen future data."
        ]},
        {"box": True, "title": "Final Decision Rule", "text": "FS ≥ τ → ANSWER (Safe to trust)\nFS < τ → ABSTAIN (Defer to human physician)"}
    ])
    
    # 9. Engineering
    add_content_slide(prs, "Experimental Setup & Engineering", [
        {"title": "Rigorous Testing Across Multiple Modalities", "bullets": [
            "Mock (Synthetic): Mathematically prove the FMR signals work perfectly.",
            "VQA-RAD: Real clinical radiology (X-rays, CTs, MRIs).",
            "PathVQA: Real pathology microscope slides.",
            "SLAKE: Multi-modal, bilingual clinical imaging."
        ]},
        {"box": True, "title": "Infrastructure Built", "text": "Executed entirely on NVIDIA T4 GPUs via Google Colab, processing thousands of complex chain-of-thought inference passes."}
    ])
    
    # 10. Live API & Dashboard
    add_content_slide(prs, "The Live Clinical API & Dashboard", [
        {"title": "Real-Time AI Auditing", "bullets": [
            "Live FastAPI Backend: We engineered a Python server hosting MedVLM-R1.",
            "Dynamic Inference: Doctors can upload unseen X-rays via an Ngrok tunnel.",
            "Interactive Dashboard: Deployed via Vercel to visualize massive datasets (AUROC charts, Risk-Coverage curves)."
        ]},
        {"box": True, "title": "Real-Time Safety Check", "text": "The backend runs all 5 consistency passes live and returns the final FMR score and Conformal Gate decision in seconds."}
    ])
    
    # 11. Conclusion
    add_content_slide(prs, "Conclusion & Impact", [
        {"title": "What We Accomplished", "bullets": [
            "Designed, engineered, and proved a complete end-to-end safety auditor.",
            "By combining Vision-Language Model inference with Conformal Prediction, we proved we can mathematically detect when an AI is hallucinating."
        ]},
        {"box": True, "title": "The Bottom Line", "text": "We provided the medical AI industry with a critical tool: a system that doesn't just ask the AI for an answer, but verifies whether the AI actually looked at the patient before speaking."}
    ])
    
    prs.save('fmr/docs/FMR_Presentation.pptx')
    print("Saved to fmr/docs/FMR_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
